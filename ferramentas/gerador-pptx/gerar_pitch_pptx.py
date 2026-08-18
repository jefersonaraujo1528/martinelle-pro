# -*- coding: utf-8 -*-
"""Gera os PowerPoints dos pitches Meta Ads Pro e YouTube Ads Pro (Consultoria MRTN).
Reconstrucao nativa em python-pptx, fiel aos HTMLs em propostas/modelos/pitch-*.html
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

SC = os.path.dirname(os.path.abspath(__file__))
DEST = "/Users/marciolinhares/Desktop/MRNT"

# paleta (igual ao CSS do pitch)
INK        = RGBColor(0x0F, 0x14, 0x1C)
MUTE       = RGBColor(0x7C, 0x86, 0x98)
MUTE2      = RGBColor(0xAA, 0xB2, 0xC0)
LINE       = RGBColor(0xE7, 0xEA, 0xF0)
SOFT       = RGBColor(0xF4, 0xF7, 0xFF)
SOFTLINE   = RGBColor(0xD4, 0xE2, 0xFF)
ACCENT     = RGBColor(0x2F, 0x6B, 0xFF)
ACCENTDEEP = RGBColor(0x1A, 0x3F, 0x99)
DARK       = RGBColor(0x0B, 0x10, 0x17)
DARK2      = RGBColor(0x16, 0x20, 0x2E)
GOLD       = RGBColor(0xE8, 0xB8, 0x4B)
GREEN      = RGBColor(0x1F, 0x9D, 0x55)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DIMBG      = RGBColor(0xFA, 0xFB, 0xFD)

SERIF = "Georgia"
SANS  = "Arial"

W, H = 13.333, 7.5

# ---------- logos ----------
LOGO_DARK  = os.path.join(SC, "logo-dark.png")
LOGO_LIGHT = os.path.join(SC, "logo-light.png")
if not os.path.exists(LOGO_LIGHT):
    im = Image.open(LOGO_DARK).convert("RGBA")
    px = im.load()
    for y in range(im.height):
        for x in range(im.width):
            r, g, b, a = px[x, y]
            px[x, y] = (255, 255, 255, a)
    im.save(LOGO_LIGHT)


def bg(slide, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def grad_dark(slide):
    bg(slide, DARK)
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2.2), Inches(-4.2), Inches(17.7), Inches(9.0))
    sh.fill.solid(); sh.fill.fore_color.rgb = DARK2
    sh.line.fill.background(); sh.shadow.inherit = False
    sh.fill.transparency = 0.35


def logo(slide, light=False):
    slide.shapes.add_picture(LOGO_LIGHT if light else LOGO_DARK,
                             Inches(0.52), Inches(0.42), height=Inches(0.30))


def txt(slide, x, y, w, h, text, size=18, color=INK, bold=False, font=SANS,
        align=PP_ALIGN.CENTER, spacing=1.15, anchor=MSO_ANCHOR.TOP, space_after=0,
        caps=False, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, linha in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        # negrito inline com **
        partes = linha.split("**")
        for k, parte in enumerate(partes):
            if parte == "":
                continue
            r = p.add_run()
            r.text = parte.upper() if caps else parte
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold or (k % 2 == 1)
            r.font.name = font
            r.font.italic = italic
    return tb


def card(slide, x, y, w, h, fill=None, line=LINE, radius=True):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        sh.adjustments[0] = 0.06
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def pill(slide, x, y, w, h, texto, fill=ACCENT, cor=WHITE, size=11):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.5
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = texto.upper()
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = cor; r.font.name = SANS
    return sh


def eyebrow(slide, texto, color=ACCENT, y=1.05):
    txt(slide, 1.0, y, W - 2.0, 0.3, texto, size=12.5, color=color, bold=True, caps=True)


def dots(slide, idx, total=9, light=False):
    d = 0.075
    gap = 0.13
    total_w = total * d + (total - 1) * (gap - d)
    x0 = (W - total_w) / 2
    for k in range(total):
        sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x0 + k * gap), Inches(H - 0.42),
                                    Inches(d), Inches(d))
        on = (k == idx)
        if light:
            cor = WHITE if on else RGBColor(0x55, 0x5F, 0x70)
        else:
            cor = INK if on else LINE
        sh.fill.solid(); sh.fill.fore_color.rgb = cor
        sh.line.fill.background(); sh.shadow.inherit = False


# =======================  SLIDES  =======================

def s_capa(slide, d):
    bg(slide, WHITE); logo(slide)
    pill(slide, (W - 2.5) / 2, 1.75, 2.5, 0.36, d["plano"])
    txt(slide, 1.2, 2.42, W - 2.4, 2.6, d["capa_h1"], size=48, color=INK, font=SERIF, spacing=1.06)
    txt(slide, 1.2, 5.35, W - 2.4, 0.4, "Consultoria MRTN · **Jeferson Martinelle** — Diretor",
        size=15, color=MUTE)
    dots(slide, 0)


def s_conceito(slide, d):
    bg(slide, WHITE); logo(slide)
    eyebrow(slide, d["c_eyebrow"])
    txt(slide, 1.0, 1.5, W - 2.0, 1.1, d["c_h1"], size=36, color=INK, font=SERIF, spacing=1.08)
    cols = d["c_cols"]
    cw, gap = 4.55, 0.55
    x0 = (W - (cw * 2 + gap)) / 2
    for i, c in enumerate(cols):
        x = x0 + i * (cw + gap)
        quente = (i == 1)
        card(slide, x, 2.95, cw, 2.35, SOFT if quente else DIMBG, SOFTLINE if quente else LINE)
        txt(slide, x + 0.35, 3.2, cw - 0.7, 0.28, c["tag"], size=11,
            color=ACCENT if quente else MUTE2, bold=True, align=PP_ALIGN.LEFT, caps=True)
        txt(slide, x + 0.35, 3.58, cw - 0.7, 0.5, c["big"], size=22, color=INK,
            font=SERIF, align=PP_ALIGN.LEFT)
        txt(slide, x + 0.35, 4.18, cw - 0.7, 1.0, c["body"], size=13, color=MUTE,
            align=PP_ALIGN.LEFT, spacing=1.35)
    txt(slide, 2.2, 5.55, W - 4.4, 0.9, d["c_lead"], size=14.5, color=MUTE, spacing=1.4)
    dots(slide, 1)


def s_numeros(slide, d):
    bg(slide, WHITE); logo(slide)
    eyebrow(slide, d["n_eyebrow"])
    txt(slide, 1.0, 1.5, W - 2.0, 1.05, d["n_h1"], size=34, color=INK, font=SERIF, spacing=1.08)
    cw, gap = 3.5, 0.42
    x0 = (W - (cw * 3 + gap * 2)) / 2
    for i, c in enumerate(d["n_cards"]):
        x = x0 + i * (cw + gap)
        card(slide, x, 3.0, cw, 1.85, SOFT, SOFTLINE)
        txt(slide, x + 0.25, 3.25, cw - 0.5, 0.42, c["num"], size=21, color=ACCENTDEEP,
            bold=True, font=SERIF)
        txt(slide, x + 0.25, 3.78, cw - 0.5, 0.95, c["txt"], size=12.5, color=INK, spacing=1.35)
    txt(slide, 1.9, 5.2, W - 3.8, 1.2, d["n_lead"], size=14, color=MUTE, spacing=1.4)
    dots(slide, 2)


def s_metodo(slide, d):
    bg(slide, WHITE); logo(slide)
    eyebrow(slide, d.get("m_eyebrow", "O método · como eu trabalho"))
    txt(slide, 1.0, 1.42, W - 2.0, 0.7, d.get("m_h1", "Um processo, não um chute."), size=32,
        color=INK, font=SERIF)
    y = 2.35
    bw = 8.6
    x = (W - bw) / 2
    for i, st in enumerate(d["metodo"]):
        star = (i == 3)
        alt = 1.05 if not star else 1.15
        card(slide, x, y, bw, alt, SOFT if star else None, SOFTLINE if star else LINE)
        txt(slide, x + 0.3, y + 0.22, 0.6, 0.4, "0%d" % (i + 1), size=20,
            color=ACCENT if star else MUTE2, font=SERIF, align=PP_ALIGN.LEFT)
        txt(slide, x + 0.95, y + 0.17, bw - 1.3, 0.3, st["t"], size=14,
            color=ACCENTDEEP if star else INK, bold=True, align=PP_ALIGN.LEFT)
        txt(slide, x + 0.95, y + 0.5, bw - 1.3, 0.6, st["d"], size=11.5, color=MUTE,
            align=PP_ALIGN.LEFT, spacing=1.3)
        y += alt + 0.16
    dots(slide, 3)


def s_deliv(slide, d):
    bg(slide, WHITE); logo(slide)
    eyebrow(slide, "O que você recebe")
    card(slide, (W - 8.2) / 2, 1.45, 8.2, 0.62, SOFT, SOFTLINE)
    txt(slide, (W - 8.0) / 2, 1.62, 8.0, 0.35,
        d.get("deliv_banner", "**Você paga só o gerenciamento.** Todo o resto é bônus por fidelidade no contrato de 12 meses."),
        size=13.5, color=ACCENTDEEP)
    cw, gap = 4.7, 0.5
    x0 = (W - (cw * 2 + gap)) / 2
    # coluna 1 - servico
    txt(slide, x0, 2.45, cw, 0.3, d.get("serv_title", "O serviço — gerenciamento mensal"), size=11,
        color=ACCENT, bold=True, align=PP_ALIGN.LEFT, caps=True)
    y = 2.9
    for i, r in enumerate(d["serv"]):
        marca = "★" if i == 0 else "—"
        cor = ACCENTDEEP if i == 0 else INK
        txt(slide, x0, y, 0.25, 0.3, marca, size=11, color=ACCENT if i == 0 else MUTE2,
            align=PP_ALIGN.LEFT)
        txt(slide, x0 + 0.3, y - 0.02, cw - 0.3, 0.55, r, size=12.5, color=cor,
            bold=(i == 0), align=PP_ALIGN.LEFT, spacing=1.25)
        y += 0.46
    # coluna 2 - bonus
    x1 = x0 + cw + gap
    txt(slide, x1, 2.45, cw, 0.3, d.get("bonus_title", "Bônus por fidelidade · 12 meses"), size=11,
        color=ACCENT, bold=True, align=PP_ALIGN.LEFT, caps=True)
    y = 2.9
    for b in d["bonus"]:
        txt(slide, x1, y, 0.25, 0.3, "★" if b.get("star") else "—", size=11,
            color=ACCENT if b.get("star") else MUTE2, align=PP_ALIGN.LEFT)
        txt(slide, x1 + 0.3, y - 0.02, cw - 1.4, 0.5, b["t"], size=12.5,
            color=ACCENTDEEP if b.get("star") else INK, bold=b.get("star", False),
            align=PP_ALIGN.LEFT, spacing=1.2)
        txt(slide, x1 + cw - 1.1, y - 0.02, 1.1, 0.3, b["v"], size=12.5, color=MUTE,
            bold=True, align=PP_ALIGN.RIGHT)
        y += 0.42
        if b.get("nota"):
            txt(slide, x1 + 0.3, y - 0.05, cw - 0.5, 0.4, b["nota"], size=10.5,
                color=MUTE2, align=PP_ALIGN.LEFT, spacing=1.25)
            y += 0.42
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x1), Inches(y + 0.05), Inches(cw), Emu(9525))
    ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background(); ln.shadow.inherit = False
    txt(slide, x1, y + 0.2, cw - 1.2, 0.3, d.get("bonus_total_label", "Total em bônus, sem custo"), size=12.5,
        color=GREEN, bold=True, align=PP_ALIGN.LEFT)
    txt(slide, x1 + cw - 1.4, y + 0.2, 1.4, 0.3, d["bonus_total"], size=12.5,
        color=MUTE2, bold=True, align=PP_ALIGN.RIGHT)
    dots(slide, 4)


def s_ancora(slide, label, valor, idx, ano=None):
    bg(slide, WHITE); logo(slide)
    y0 = 1.75 if ano else 2.15
    txt(slide, 2.6, y0, W - 5.2, 0.9, label, size=15, color=MUTE, bold=True, spacing=1.4, caps=True)
    txt(slide, 1.0, y0 + 1.0, W - 2.0, 1.5, valor, size=88 if ano else 96, color=INK, font=SERIF, spacing=1.0)
    txt(slide, 1.0, y0 + 2.55, W - 2.0, 0.35, "por mês", size=15, color=MUTE)
    if ano:
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.4), Inches(5.05), Inches(4.5), Emu(9525))
        ln.fill.solid(); ln.fill.fore_color.rgb = LINE; ln.line.fill.background(); ln.shadow.inherit = False
        txt(slide, 1.0, 5.28, W - 2.0, 0.3, "que no ano vira", size=14, color=MUTE)
        txt(slide, 1.0, 5.62, W - 2.0, 0.75, ano["v"], size=40, color=INK, font=SERIF)
        txt(slide, 1.0, 6.42, W - 2.0, 0.35, ano["extra"], size=14, color=MUTE)
    dots(slide, idx)


def s_canais(slide, d, idx=6):
    bg(slide, WHITE); logo(slide)
    eyebrow(slide, "Escolha o seu canal — e leve o bônus dele")
    txt(slide, 1.0, 1.45, W - 2.0, 0.7, "Um canal. Um presente.", size=32, color=INK, font=SERIF)
    cw, gap = 4.9, 0.55
    x0 = (W - (cw * 2 + gap)) / 2
    for i, c in enumerate(d["canais"]):
        x = x0 + i * (cw + gap)
        card(slide, x, 2.35, cw, 3.7, None, LINE)
        card(slide, x, 4.25, cw, 1.8, SOFT, SOFTLINE)
        txt(slide, x + 0.35, 2.58, cw - 0.7, 0.28, c["kick"], size=10.5, color=MUTE2, bold=True,
            align=PP_ALIGN.LEFT, caps=True)
        txt(slide, x + 0.35, 2.95, cw - 0.7, 0.7, c["nome"], size=20, color=INK, font=SERIF,
            align=PP_ALIGN.LEFT, spacing=1.12)
        txt(slide, x + 0.35, 3.72, cw - 0.7, 0.5, c["oq"], size=12, color=MUTE,
            align=PP_ALIGN.LEFT, spacing=1.35)
        txt(slide, x + 0.35, 4.45, cw - 0.7, 0.26, "★ SEU BÔNUS", size=10, color=ACCENT, bold=True,
            align=PP_ALIGN.LEFT)
        txt(slide, x + 0.35, 4.76, cw - 0.7, 0.3, c["gn"], size=13.5, color=ACCENTDEEP, bold=True,
            align=PP_ALIGN.LEFT)
        txt(slide, x + 0.35, 5.1, cw - 0.7, 0.55, c["gd"], size=11, color=MUTE,
            align=PP_ALIGN.LEFT, spacing=1.3)
        pill(slide, x + 0.35, 5.68, 2.1, 0.3, c["gp"], fill=WHITE, cor=GREEN, size=11)
    txt(slide, 1.0, 6.35, W - 2.0, 0.4,
        "O bônus é da **decisão em 24 horas** — depois disso o plano continua igual, o presente não.",
        size=13.5, color=RGBColor(0x8A, 0x5D, 0x00))
    dots(slide, idx)


def s_vp(slide, d):
    grad_dark(slide); logo(slide, light=True)
    eyebrow(slide, "O que importa de verdade", color=GOLD)
    # lado valor
    txt(slide, 1.5, 2.3, 4.2, 0.35, d.get("vp_label_esq", "Quanto vale"), size=13, color=GOLD, bold=True, caps=True)
    txt(slide, 1.5, 2.95, 4.2, 0.35, d["vp_l1"], size=15, color=RGBColor(0xCF, 0xD6, 0xE2))
    txt(slide, 1.5, 3.4, 4.2, 0.35, d["vp_l2"], size=15, color=RGBColor(0xCF, 0xD6, 0xE2))
    txt(slide, 1.5, 3.95, 4.2, 1.1, "muito", size=54, color=GOLD, font=SERIF)
    txt(slide, 1.5, 5.15, 4.2, 0.3, "valor real entregue", size=12.5, color=RGBColor(0x9A, 0xA5, 0xB5))
    # x
    txt(slide, (W - 1.0) / 2, 3.75, 1.0, 0.8, "×", size=40, color=RGBColor(0x6B, 0x76, 0x88), font=SERIF)
    # lado preco
    txt(slide, W - 5.7, 2.3, 4.2, 0.35, d.get("vp_label_dir", "Quanto você paga"), size=13, color=RGBColor(0xD8, 0xDE, 0xE8),
        bold=True, caps=True)
    txt(slide, W - 5.7, 3.25, 4.2, 1.5, d["vp_mes"], size=58, color=WHITE, font=SERIF)
    txt(slide, W - 5.7, 4.85, 4.2, 0.4, d["vp_sub"], size=13, color=RGBColor(0x9A, 0xA5, 0xB5))
    dots(slide, 7, light=True)


def s_final(slide, d):
    bg(slide, ACCENT); logo(slide, light=True)
    # bloco preco (esquerda)
    txt(slide, 0.95, 1.85, 5.2, 0.3, d.get("f_lbl", "O nosso valor"), size=12.5,
        color=RGBColor(0xCF, 0xDD, 0xFF), bold=True, caps=True)
    txt(slide, 0.95, 2.3, 5.2, 0.35, d["f_anc"], size=13, color=RGBColor(0xC2, 0xD4, 0xFF))
    txt(slide, 0.95, 2.75, 5.2, 1.6, d["f_val"], size=76, color=WHITE, font=SERIF)
    txt(slide, 0.95, 4.5, 5.2, 0.4, d["f_per"], size=13.5, color=RGBColor(0xE4, 0xEC, 0xFF))
    txt(slide, 0.95, 5.1, 5.2, 0.9, d["f_close"], size=13, color=WHITE, spacing=1.4, italic=True)
    # stack (direita)
    bx, bw = 6.9, 5.5
    card(slide, bx, 1.75, bw, 4.6, RGBColor(0x4B, 0x7F, 0xFF), RGBColor(0x8F, 0xB2, 0xFF))
    y = 2.05
    for bloco in d["f_stack"]:
        txt(slide, bx + 0.42, y, bw - 0.84, 0.3, bloco["t"], size=10.5,
            color=RGBColor(0xD3, 0xE0, 0xFF), bold=True, align=PP_ALIGN.LEFT, caps=True)
        y += 0.38
        for it in bloco["itens"]:
            txt(slide, bx + 0.42, y, 0.25, 0.28, "★" if it.get("star") else "—", size=10.5,
                color=RGBColor(0xCF, 0xDD, 0xFF), align=PP_ALIGN.LEFT)
            txt(slide, bx + 0.72, y - 0.03, bw - 2.2, 0.5, it["t"], size=12, color=WHITE,
                align=PP_ALIGN.LEFT, spacing=1.2)
            if it.get("v"):
                txt(slide, bx + bw - 1.5, y - 0.03, 1.1, 0.3, it["v"], size=11.5,
                    color=RGBColor(0xC2, 0xD4, 0xFF), align=PP_ALIGN.RIGHT)
            y += 0.42 if len(it["t"]) < 52 else 0.62
        y += 0.14
    dots(slide, 8, light=True)


# =======================  CONTEUDO  =======================

META = dict(
    arquivo="Pitch Meta Ads Pro - Consultoria MRTN.pptx",
    titulo="Pitch Meta Ads Pro — Consultoria MRTN",
    plano="Meta Ads Pro",
    capa_h1="Fazer o seu nome chegar onde\na sua agenda ainda não chega.",
    c_eyebrow="A diferença que muda tudo",
    c_h1="Esperar ser procurado é um plano ruim.",
    c_cols=[
        dict(tag="Sem anúncio", big="Você espera.",
             body="A agenda depende de indicação, de quem passou na porta, de quem ouviu falar. Um mês cheio, o outro vazio — e você nunca sabe explicar o porquê."),
        dict(tag="Com Meta Ads", big="Você decide.",
             body="Você escolhe quem vê, em qual bairro, com que idade, quantas vezes e por quanto tempo. A agenda para de depender da sorte e passa a depender de um processo."),
    ],
    c_lead="Essa é a diferença entre ter um consultório que reage ao mês e ter um negócio que cresce por decisão.",
    n_eyebrow="Não é opinião — é o mercado",
    n_h1="O seu paciente já está lá. Todo dia.",
    n_cards=[
        dict(num="134 milhões", txt="de brasileiros no Instagram — 3º maior mercado do mundo, com alcance de 69% da população"),
        dict(num="73%", txt="dos usuários já contrataram um serviço que descobriram dentro da plataforma"),
        dict(num="46% em 7 dias", txt="de quem descobre algo lá, 89% pesquisam mais, 67% visitam o site e 46% contratam em uma semana"),
    ],
    n_lead="Não é sobre estar na moda. É sobre estar presente no lugar onde a decisão nasce — e continuar aparecendo até a pessoa precisar de você.",
    metodo=[
        dict(t="Estratégia & Setup", d="Construo os públicos, os interesses e as segmentações do seu consultório — e subo até 5 campanhas com 5 conjuntos e no mínimo 3 anúncios cada. Testo de verdade em vez de apostar em um só criativo."),
        dict(t="Rastreamento — o que quase ninguém instala", d="Pixel do Meta configurado, mais a sua página de bio com o Pixel dentro. Cada pessoa que passa por ali vira público de remarketing: quem olhou e não chamou volta a ver você."),
        dict(t="Presença que vira reconhecimento", d="Gerencio e monitoro todo dia — custo por resultado, alcance, criativo que cansou. Ninguém fecha com quem viu uma vez: fecha com quem já parece familiar."),
        dict(t="★ Consultoria mensal — o coração do trabalho", d="Todo mês, sentamos por vídeo e analisamos juntos os números e os gargalos do seu funil. É aqui que o crescimento vira decisão — e você nunca fica na mão."),
    ],
    serv=["Consultoria mensal por vídeo comigo",
          "Gerenciamento e monitoramento diário das campanhas",
          "Públicos ilimitados e campanhas de remarketing",
          "Anúncios otimizados dentro das normas do Facebook e Instagram",
          "Relatórios mensais e semanais + otimização contínua"],
    bonus=[dict(t="Setup completo das campanhas", v="R$ 2.000"),
           dict(t="Website Smart Links — a página da sua bio", v="R$ 1.500", star=True,
                nota="WhatsApp, localização e agendamento num só lugar — com Pixel instalado.")],
    bonus_total="R$ 3.500",
    a1_label="O que o mercado cobra\npor uma gestão especializada de Meta Ads",
    a1_valor="R$ 2.500",
    a2_label="A taxa de tabela\nda Consultoria MRTN",
    a2_valor="R$ 1.500",
    vp_l1="Trabalho: **R$ 1.500/mês**",
    vp_l2="+ bônus: **R$ 3.500**",
    vp_mes="R$ 1.200",
    vp_sub="por mês · o equivalente a R$ 40 por dia",
    f_anc="Mercado R$ 2.500 · tabela R$ 1.500",
    f_val="R$ 1.200",
    f_per="por mês · o equivalente a R$ 40 por dia\nno Pix ou boleto, contrato de 12 meses",
    f_close="Você paga só o gerenciamento. O resto é meu presente pela sua fidelidade.",
    f_stack=[
        dict(t="Você paga só o gerenciamento", itens=[
            dict(t="Consultoria mensal por vídeo comigo", star=True),
            dict(t="Monitoramento diário, públicos ilimitados, remarketing"),
            dict(t="Relatórios mensais e semanais + otimização")]),
        dict(t="Bônus por fidelidade · 12 meses — R$ 3.500 grátis", itens=[
            dict(t="Setup completo das campanhas", v="R$ 2.000"),
            dict(t="Website Smart Links com Pixel", v="R$ 1.500", star=True)]),
    ],
)

YT = dict(
    arquivo="Pitch YouTube Ads Pro - Consultoria MRTN.pptx",
    titulo="Pitch YouTube Ads Pro — Consultoria MRTN",
    plano="YouTube Ads Pro",
    capa_h1="Ocupar a maior tela\nda casa do seu paciente.",
    c_eyebrow="A diferença que muda tudo",
    c_h1="Currículo informa. Rosto convence.",
    c_cols=[
        dict(tag="Texto e imagem", big="Você é uma informação.",
             body="Nome, especialidade, endereço. A pessoa lê, compara com outros três e decide pelo que parecer mais perto ou mais barato. Você virou item de lista."),
        dict(tag="Vídeo", big="Você é uma pessoa.",
             body="Ela ouve como você explica, vê como você olha, sente o seu jeito. Chega no consultório já confiando — e quem confia não pergunta o preço primeiro."),
    ],
    c_lead="Saúde não se escolhe pelo menor preço. Se escolhe pela maior confiança. E não existe formato que construa confiança como alguém ver o seu rosto e ouvir a sua voz.",
    n_eyebrow="Não é opinião — é o mercado",
    n_h1="O YouTube saiu do celular e voltou pra sala.",
    n_cards=[
        dict(num="144 milhões", txt="de brasileiros no YouTube — 3º maior público do mundo, atrás só de Índia e Estados Unidos"),
        dict(num="80 milhões", txt="assistem pela TV da sala. A TV conectada ultrapassou o celular como tela principal"),
        dict(num="41% → 53%", txt="foi o salto da audiência do YouTube na TV conectada entre adultos, em três anos"),
    ],
    n_lead="Mais adultos assistem YouTube numa semana comum do que cada uma das cinco maiores emissoras de TV aberta. Só que na TV você não escolhe quem vê — aqui, você escolhe a cidade, o bairro, a idade e o interesse. É comercial de televisão com pontaria.",
    metodo=[
        dict(t="Estratégia & roteiro", d="Defino o que cada vídeo precisa dizer e em qual formato — in-stream, bumper ou discovery. Os vídeos são seus; a orientação de roteiro e de formato é minha."),
        dict(t="★ O Google Ads Pro inteiro vem junto", d="Você não leva só o vídeo: leva o pacote completo de busca — campanhas no Google, Google Meu Negócio otimizado e landing page de alta conversão. O vídeo planta o nome; a busca colhe quando a pessoa vai procurar."),
        dict(t="Campanhas no ar", d="Gerencio diariamente — canais, tópicos, públicos, CPV, retenção. Eu vejo até em que segundo a pessoa para de assistir, e corrijo a partir disso."),
        dict(t="★ Consultoria mensal — o coração do trabalho", d="Todo mês, sentamos por vídeo e analisamos juntos os números e os gargalos do seu funil. É aqui que o crescimento vira decisão — e você nunca fica na mão."),
    ],
    serv=["Consultoria mensal por vídeo comigo",
          "Gerenciamento e monitoramento diário das campanhas",
          "Gerenciamento diário das campanhas de vídeo e de busca",
          "CPV, CTR, retenção e conversões monitorados de perto",
          "Relatórios mensais e semanais + otimização contínua"],
    bonus=[dict(t="Setup das campanhas de vídeo", v="R$ 6.000"),
           dict(t="O pacote Google Ads Pro completo", v="R$ 5.000", star=True,
                nota="Campanhas de busca, Google Meu Negócio e landing page — o pacote inteiro, dentro do seu.")],
    bonus_total="R$ 11.000",
    a1_label="O que o mercado cobra\npor uma gestão especializada de YouTube Ads",
    a1_valor="R$ 5.000",
    a2_label="A taxa de tabela\nda Consultoria MRTN",
    a2_valor="R$ 4.500",
    vp_l1="Trabalho: **R$ 4.500/mês**",
    vp_l2="+ bônus: **R$ 11.000**",
    vp_mes="R$ 3.500",
    vp_sub="por mês · o equivalente a R$ 117 por dia",
    f_anc="Mercado R$ 5.000 · tabela R$ 4.500",
    f_val="R$ 3.500",
    f_per="por mês · o equivalente a R$ 117 por dia\nno Pix ou boleto, contrato de 12 meses",
    f_close="Dois pacotes num só: o YouTube constrói o seu nome, o Google Ads Pro colhe quem procura. Você paga só o gerenciamento.",
    f_stack=[
        dict(t="Você paga só o gerenciamento", itens=[
            dict(t="Consultoria mensal por vídeo comigo", star=True),
            dict(t="Gerenciamento diário das campanhas de vídeo e de busca"),
            dict(t="Relatórios mensais e semanais + otimização")]),
        dict(t="Bônus por fidelidade · 12 meses — R$ 11.000 grátis", itens=[
            dict(t="Setup das campanhas de vídeo", v="R$ 6.000"),
            dict(t="O pacote Google Ads Pro completo", v="R$ 5.000", star=True)]),
    ],
)


GOOGLE = dict(
    arquivo="Pitch Google Ads Pro - Consultoria MRTN.pptx",
    titulo="Pitch Google Ads Pro — Consultoria MRTN",
    plano="Google Ads Pro",
    capa_h1="Fazer o paciente certo encontrar\nvocê na hora que ele procura.",
    c_eyebrow="A diferença que muda tudo",
    c_h1="Atenção não é intenção.",
    c_cols=[
        dict(tag="Redes sociais · atenção", big="Você interrompe.",
             body="A pessoa está distraída, vendo outra coisa. Você aparece no meio e torce pra ela parar. Ela não estava procurando por você."),
        dict(tag="Google · intenção", big="O paciente te procura.",
             body="No exato momento em que alguém digita a sua especialidade na sua cidade, ela já quer resolver. Você não empurra — você aparece na hora certa."),
    ],
    c_lead="É por isso que o Google é o canal do médico: quem busca já decidiu que precisa. Falta só te encontrar.",
    n_eyebrow="Não é opinião — é o mercado",
    n_h1="Cada busca é uma oportunidade de agenda.",
    n_cards=[
        dict(num="92%", txt="das buscas feitas no Brasil acontecem dentro do Google. É lá que a decisão começa"),
        dict(num="85%", txt="dos pacientes pesquisam na internet antes de marcar uma consulta"),
        dict(num="45%", txt="da população brasileira busca informação de saúde online — todo mês, na sua cidade"),
    ],
    n_lead="O paciente precisa, pesquisa agora, encontra a sua especialidade e agenda. Meu trabalho é fazer você aparecer nesse caminho — e ser escolhido no lugar do concorrente.",
    metodo=[
        dict(t="Estratégia & Setup", d="Pesquiso as palavras-chave que seus pacientes realmente usam, estudo seus concorrentes e estruturo as campanhas do zero."),
        dict(t="Presença local", d="Configuro e otimizo seu Google Meu Negócio toda semana, com ferramentas avançadas para você dominar o seu bairro na busca."),
        dict(t="Campanhas no ar", d="Gerencio e otimizo diariamente — anúncios, lances, palavras-chave — para cada real investido render mais."),
        dict(t="★ Consultoria mensal — o coração do trabalho", d="Todo mês, sentamos por vídeo e analisamos juntos os números e os gargalos do seu funil. É aqui que o crescimento vira decisão — e você nunca fica na mão."),
    ],
    serv=["Consultoria mensal por vídeo comigo",
          "Gerenciamento e monitoramento diário das campanhas",
          "Palavras-chave ilimitadas, remarketing e geolocalização",
          "Anúncios otimizados para mais cliques e conversões",
          "Relatórios mensais e semanais + otimização contínua"],
    bonus=[dict(t="Setup completo das campanhas", v="R$ 3.000"),
           dict(t="Google Meu Negócio + otimização semanal", v="R$ 800", star=True,
                nota="Dominar o seu bairro e a sua região na busca local."),
           dict(t="Landing page de alta conversão", v="R$ 1.200")],
    bonus_total="R$ 5.000",
    a1_label="O que o mercado cobra\npor uma publicidade especializada no Google",
    a1_valor="R$ 3.000",
    a2_label="O que esse trabalho\nrealmente vale — com tudo incluído",
    a2_valor="R$ 1.800",
    vp_l1="Trabalho: **R$ 1.800/mês**",
    vp_l2="+ bônus: **R$ 5.000**",
    vp_mes="R$ 1.200",
    vp_sub="por mês · o equivalente a R$ 40 por dia",
    f_anc="Mercado R$ 3.000 · vale R$ 1.800",
    f_val="R$ 1.200",
    f_per="por mês · o equivalente a R$ 40 por dia\nno Pix ou boleto, contrato de 12 meses",
    f_close="Você paga só o gerenciamento. O resto é meu presente pela sua fidelidade.",
    f_stack=[
        dict(t="Você paga só o gerenciamento", itens=[
            dict(t="Consultoria mensal por vídeo comigo", star=True),
            dict(t="Monitoramento diário, palavras-chave ilimitadas, remarketing"),
            dict(t="Relatórios mensais e semanais + otimização")]),
        dict(t="Bônus por fidelidade · 12 meses — R$ 5.000 grátis", itens=[
            dict(t="Setup completo das campanhas", v="R$ 3.000"),
            dict(t="Google Meu Negócio + otimização semanal", v="R$ 800", star=True),
            dict(t="Landing page de alta conversão", v="R$ 1.200")]),
    ],
)


LOCAL = dict(
    arquivo="Apresentacao Negocios Locais - Consultoria MRTN.pptx",
    titulo="Apresentação Negócios Locais — Consultoria MRTN",
    plano="Negócios Locais",
    capa_h1="Todo dia alguém procura o que você vende.\nA questão é quem aparece.",
    c_eyebrow="A diferença que muda tudo",
    c_h1="Quem não aparece, não é escolhido.",
    c_cols=[
        dict(tag="Só indicação", big="Você espera.",
             body="Um mês cheio, o outro parado — e sem saber explicar o porquê. O movimento depende de quem lembrou de você, de quem passou na porta, da sorte da semana."),
        dict(tag="Com publicidade no Google e no Instagram", big="Você escolhe.",
             body="Escolhe quem vê o seu negócio: qual bairro, qual idade, qual interesse — e quantas vezes. Deixa de torcer e passa a ter um processo que você controla."),
    ],
    c_lead="Não é sobre aparecer para todo mundo. É sobre aparecer para quem já está procurando — no bairro certo e na hora certa.",
    n_eyebrow="Não é opinião — é o mercado",
    n_h1="O seu cliente já está lá. Hoje.",
    n_cards=[
        dict(num="92%", txt="das buscas feitas no Brasil acontecem dentro do Google. É lá que a procura começa"),
        dict(num="134 milhões", txt="de brasileiros no Instagram — alcance de 69% da população do país"),
        dict(num="73%", txt="dos usuários já compraram algo que descobriram na plataforma, sem estar procurando"),
    ],
    m_eyebrow="Como começa · o processo",
    m_h1="Nada vai ao ar sem você aprovar.",
    deliv_banner="**Você paga só a gestão da publicidade.** A construção do funil e o bônus entram sem custo.",
    serv_title="Todo mês, o serviço",
    bonus_title="Sem custo, no fechamento",
    bonus_total_label="Total sem custo",
    f_lbl="Plano Start · 1 canal",
    n_lead="São os dois lugares onde a decisão de compra nasce: o que a pessoa procura e o que ela descobre. Estar nos dois é o que separa o negócio conhecido do negócio lembrado.",
    metodo=[
        dict(t="Briefing e aprovação", d="Entendo o seu negócio, o seu público e o que você quer vender mais. Você confirma que eu entendi certo antes de qualquer coisa começar."),
        dict(t="Construção do funil", d="Monto as campanhas, os públicos e o caminho que leva a pessoa do anúncio até o seu WhatsApp. Você vê pronto e aprova."),
        dict(t="Campanhas no ar", d="Acompanho todo dia — o que gasta sem retornar sai, o que traz cliente ganha verba. Ajuste é feito na hora, não no fim do mês."),
        dict(t="★ WhatsApp direto comigo, e relatório toda semana", d="Você fala comigo pelo WhatsApp de segunda a sexta, das 9h às 18h — e recebe relatório semanal e mensal mostrando o que foi investido e o que voltou."),
    ],
    serv=["WhatsApp direto comigo — 2ª a 6ª, das 9h às 18h",
          "Criação e gestão dos seus anúncios no Google ou no Instagram",
          "Relatórios semanais e mensais do que foi investido e do que voltou",
          "Acompanhamento e ajuste das campanhas todos os dias"],
    bonus=[dict(t="Construção e configuração do funil", v="R$ 1.500"),
           dict(t="Bônus de decisão em 24 horas", v="até R$ 1.500", star=True,
                nota="Você escolhe o canal — e o bônus vem junto com ele.")],
    bonus_total="até R$ 3.000",
    a1_label="O que uma agência chega a cobrar\nsó pela gestão de um canal",
    a1_valor="R$ 1.500",
    a1_ano=dict(v="R$ 18.000", extra="e com a construção do funil e o bônus, R$ 21.000"),
    canais=[
        dict(kick="Opção 1 · Google", nome="Aparecer para\nquem já procura.",
             oq="Seus anúncios no topo do Google quando alguém da sua cidade busca exatamente o que você vende.",
             gn="Uma página feita para vender",
             gd="Landing page de alta conversão, com botão direto para o seu WhatsApp.",
             gp="R$ 1.500  →  grátis"),
        dict(kick="Opção 2 · Instagram e Facebook", nome="Aparecer para quem\nainda não te conhece.",
             oq="Seus anúncios no feed e nos stories de quem mora perto de você — inclusive de quem nunca ouviu falar do seu negócio.",
             gn="Consultoria de criativos que convertem",
             gd="Sessão especializada no seu segmento: o que anunciar, como falar e o que só queima verba.",
             gp="R$ 1.000  →  grátis"),
    ],
    vp_label_esq="O que isso custaria separado",
    vp_l1="Gestão 12 meses: **R$ 18.000**",
    vp_l2="+ funil e bônus: **R$ 3.000**",
    vp_risco="R$ 21.000",
    vp_sub_esq="Você não paga isso.",
    vp_label_dir="Você paga",
    vp_mes="12x R$ 650",
    vp_sub="o equivalente a R$ 21,66 por dia",
    f_anc="Mercado R$ 1.500/mês · R$ 21.000 no ano",
    f_val="12x R$ 650",
    f_per="no cartão · ou R$ 650 por mês no Pix\no equivalente a R$ 21,66 por dia",
    f_close="Decidindo em 24 horas, o bônus do seu canal é seu. Depois disso o plano continua igual — o presente não.",
    f_stack=[
        dict(t="Todo mês", itens=[
            dict(t="WhatsApp direto comigo, 2ª a 6ª, 9h–18h", star=True),
            dict(t="Criação e gestão dos seus anúncios"),
            dict(t="Relatórios semanais e mensais")]),
        dict(t="Sem custo — até R$ 3.000", itens=[
            dict(t="Construção do funil", v="R$ 1.500"),
            dict(t="Bônus do seu canal", v="até R$ 1.500", star=True)]),
        dict(t="Quer os dois canais?", itens=[
            dict(t="Plano Plus — 12x R$ 1.300, o 2º canal pelo mesmo preço")]),
    ],
)


def build(d):
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    branco = prs.slide_layouts[6]
    s_capa(prs.slides.add_slide(branco), d)
    s_conceito(prs.slides.add_slide(branco), d)
    s_numeros(prs.slides.add_slide(branco), d)
    s_metodo(prs.slides.add_slide(branco), d)
    s_deliv(prs.slides.add_slide(branco), d)
    s_ancora(prs.slides.add_slide(branco), d["a1_label"], d["a1_valor"], 5, d.get("a1_ano"))
    if d.get("canais"):
        s_canais(prs.slides.add_slide(branco), d, 6)
    else:
        s_ancora(prs.slides.add_slide(branco), d["a2_label"], d["a2_valor"], 6)
    s_vp(prs.slides.add_slide(branco), d)
    s_final(prs.slides.add_slide(branco), d)
    prs.core_properties.title = d["titulo"]
    prs.core_properties.author = "Consultoria MRTN — Jeferson Martinelle"
    caminho = os.path.join(DEST, d["arquivo"])
    prs.save(caminho)
    return caminho, len(prs.slides.__iter__.__self__._sldIdLst)


for deck in (GOOGLE, META, YT, LOCAL):
    caminho, n = build(deck)
    print("OK:", caminho, "|", n, "slides")
